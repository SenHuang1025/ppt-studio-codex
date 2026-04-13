from __future__ import annotations

import csv
import json
from pathlib import Path

import fitz
from docx import Document
from openpyxl import Workbook
from PIL import Image
from pptx import Presentation
from pptx.util import Inches


def write_text_file(path: Path, text: str, *, encoding: str = "utf-8") -> Path:
    path.write_text(text, encoding=encoding)
    return path


def write_json_file(path: Path, payload: object) -> Path:
    path.write_text(json.dumps(payload, ensure_ascii=False), encoding="utf-8")
    return path


def write_csv_file(path: Path) -> Path:
    with path.open("w", encoding="utf-8-sig", newline="") as file_handle:
        writer = csv.writer(file_handle)
        writer.writerow(["Month", "Revenue", "Cost"])
        writer.writerow(["Jan", 100, 70])
        writer.writerow(["Feb", 120, 80])
        writer.writerow(["Mar", 140, 90])
    return path


def write_excel_file(path: Path) -> Path:
    workbook = Workbook()
    sales_sheet = workbook.active
    sales_sheet.title = "Sales"
    sales_sheet.append(["Month", "Revenue", "Cost"])
    sales_sheet.append(["Jan", 100, 70])
    sales_sheet.append(["Feb", 120, 80])
    sales_sheet.append(["Mar", 140, 90])

    summary_sheet = workbook.create_sheet("Summary")
    summary_sheet.append(["Metric", "Value"])
    summary_sheet.append(["Total Revenue", 360])
    summary_sheet.append(["Total Cost", 240])

    workbook.save(path)
    return path


def write_image_file(
    path: Path,
    *,
    format_name: str = "PNG",
    size: tuple[int, int] = (320, 240),
    color: tuple[int, int, int] = (240, 210, 160),
) -> Path:
    image = Image.new("RGB", size, color=color)
    image.save(path, format=format_name)
    return path


def write_docx_file(path: Path) -> Path:
    document = Document()
    document.add_heading("Quarterly Summary", level=1)
    document.add_paragraph("Revenue increased by 18 percent across the quarter.")
    table = document.add_table(rows=3, cols=2)
    table.rows[0].cells[0].text = "Metric"
    table.rows[0].cells[1].text = "Value"
    table.rows[1].cells[0].text = "Revenue"
    table.rows[1].cells[1].text = "360"
    table.rows[2].cells[0].text = "Cost"
    table.rows[2].cells[1].text = "240"

    image_path = path.with_suffix(".png")
    write_image_file(image_path)
    document.add_picture(str(image_path))
    document.save(path)
    return path


def write_pdf_file(path: Path) -> Path:
    document = fitz.open()
    page = document.new_page()
    page.insert_text((72, 72), "Quarterly Summary")
    page.insert_text((72, 100), "Revenue increased by 18 percent.")
    document.save(path)
    document.close()
    return path


def write_pptx_file(path: Path) -> Path:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Quarterly Summary"
    slide.placeholders[1].text = "Revenue increased by 18 percent.\nCost remained stable."

    try:
        slide.notes_slide.notes_text_frame.text = "Speaker note about the growth trend."
    except Exception:
        pass

    image_path = path.with_suffix(".png")
    write_image_file(image_path)
    slide.shapes.add_picture(str(image_path), Inches(5), Inches(1.5), width=Inches(2.2), height=Inches(1.6))
    presentation.save(path)
    return path
