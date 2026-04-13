from __future__ import annotations

import asyncio

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from app.parsers.base import BaseParser, ParseResult
from app.parsers.utils import normalize_file_type, truncate_text, unique_strings


class PPTXParser(BaseParser):
    async def parse(self, file_path: str) -> ParseResult:
        return await asyncio.to_thread(self._parse_sync, file_path)

    def _parse_sync(self, file_path: str) -> ParseResult:
        presentation = Presentation(file_path)
        slides: list[dict[str, object]] = []
        text_sections: list[str] = []
        key_points: list[str] = []

        for slide_number, slide in enumerate(presentation.slides, start=1):
            title_shape = slide.shapes.title
            title = title_shape.text.strip() if title_shape and title_shape.text else ""
            text_blocks: list[str] = []
            image_count = 0
            table_count = 0
            chart_count = 0

            for shape in slide.shapes:
                text = getattr(shape, "text", "")
                if isinstance(text, str) and text.strip():
                    text_blocks.append(text.strip())
                if getattr(shape, "has_table", False):
                    table_count += 1
                if getattr(shape, "has_chart", False):
                    chart_count += 1
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image_count += 1

            notes = self._extract_notes(slide)
            unique_text_blocks = unique_strings(text_blocks, max_items=10, max_length=240)
            slides.append(
                {
                    "slide_number": slide_number,
                    "title": title,
                    "text_blocks": unique_text_blocks,
                    "text_block_count": len(unique_text_blocks),
                    "element_count": len(slide.shapes),
                    "has_title": bool(title),
                    "image_count": image_count,
                    "table_count": table_count,
                    "chart_count": chart_count,
                    "notes": truncate_text(notes, 240) if notes else "",
                    "has_notes": bool(notes),
                }
            )

            slide_lines = [f"Slide {slide_number}: {title}" if title else f"Slide {slide_number}"]
            slide_lines.extend(unique_text_blocks)
            if notes:
                slide_lines.append(f"Notes: {notes}")
            text_sections.append("\n".join(slide_lines))

            if title:
                key_points.append(f"Slide {slide_number}: {title}")
            elif unique_text_blocks:
                key_points.append(f"Slide {slide_number}: {unique_text_blocks[0]}")

        return ParseResult(
            file_type=normalize_file_type(file_path),
            summary=f"PPTX presentation with {len(slides)} slides.",
            text_content="\n\n".join(text_sections),
            structured_data={
                "slide_count": len(slides),
                "slides": slides,
            },
            key_points=unique_strings(key_points, max_items=5),
        )

    def _extract_notes(self, slide) -> str:
        try:
            notes_slide = slide.notes_slide
            notes_text_frame = getattr(notes_slide, "notes_text_frame", None)
            if notes_text_frame is None:
                return ""
            return notes_text_frame.text.strip()
        except Exception:
            return ""
